#!/usr/bin/env python3
"""毕业答辩PPT生成脚本 - 基于SaaS模式的智慧健身管理系统"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ========== 颜色方案 ==========
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)      # 主色-蓝
PRIMARY_DARK = RGBColor(0x0D, 0x33, 0x8A)  # 深蓝
ACCENT = RGBColor(0xE8, 0x6C, 0x00)        # 橙色强调
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
BG_BLUE = RGBColor(0xE8, 0xEE, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== 工具函数 ==========
def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=BLACK, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=BLACK, spacing=Pt(8), bold_prefix=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        if bold_prefix and '：' in item:
            prefix, rest = item.split('：', 1)
            run1 = p.add_run()
            run1.text = prefix + '：'
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = PRIMARY
            run1.font.name = '微软雅黑'
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = '微软雅黑'
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = '微软雅黑'
    return txBox

def add_card(slide, left, top, width, height, title, content_items,
             title_color=WHITE, bg_color=PRIMARY, content_color=BLACK,
             content_bg=WHITE, font_size=14):
    # 卡片标题
    header = add_shape(slide, left, top, width, Inches(0.6), bg_color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3),
                 Inches(0.5), title, font_size=16, bold=True, color=title_color,
                 alignment=PP_ALIGN.CENTER)
    # 卡片内容
    body = add_shape(slide, left, top + Inches(0.6), width, height - Inches(0.6), content_bg)
    body.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    add_bullet_list(slide, left + Inches(0.15), top + Inches(0.7),
                    width - Inches(0.3), height - Inches(0.8),
                    content_items, font_size=font_size, color=content_color,
                    bold_prefix=False)
    return header, body

def add_section_header(slide, title, subtitle=None):
    """添加统一的页面标题栏"""
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.2), PRIMARY)
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.6),
                 title, font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.75), Inches(10), Inches(0.35),
                     subtitle, font_size=16, color=RGBColor(0xBB, 0xCC, 0xEE))
    # 底部装饰线
    add_shape(slide, Inches(0), Inches(1.2), prs.slide_width, Inches(0.04), ACCENT)

def add_page_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f'{num}/{total}', font_size=12, color=GRAY, alignment=PP_ALIGN.RIGHT)

TOTAL_PAGES = 10

# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

# 大背景色块
add_shape(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, PRIMARY_DARK)

# 装饰条
add_shape(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06), ACCENT)
add_shape(slide, Inches(0), Inches(4.8), Inches(13.333), Inches(0.06), ACCENT)

# 学校信息
add_text_box(slide, Inches(0), Inches(0.8), prs.slide_width, Inches(0.6),
             '信息科学与工程学院', font_size=22, color=RGBColor(0xBB, 0xCC, 0xEE),
             alignment=PP_ALIGN.CENTER)

# 标题
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
             '基于SaaS模式的智慧健身管理系统', font_size=40, bold=True,
             color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.7), Inches(10.3), Inches(0.4),
             '的设计与实现', font_size=24, color=RGBColor(0xBB, 0xCC, 0xEE),
             alignment=PP_ALIGN.CENTER)

# 学生信息
info_lines = [
    '学生姓名：陈斌          学号：202203010223',
    '专业班级：计算机科学与技术 计算机2202',
    '学校指导教师：王京文      企业指导教师：文子穰',
    '学习企业：数字马力信息技术有限公司',
]
y = Inches(3.4)
for line in info_lines:
    add_text_box(slide, Inches(3.0), y, Inches(7.3), Inches(0.45),
                 line, font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
    y += Inches(0.35)

add_text_box(slide, Inches(0), Inches(5.5), prs.slide_width, Inches(0.5),
             '2025年5月', font_size=20, color=RGBColor(0xBB, 0xCC, 0xEE),
             alignment=PP_ALIGN.CENTER)

# 底部
add_text_box(slide, Inches(0), Inches(6.5), prs.slide_width, Inches(0.5),
             '毕业设计答辩', font_size=16, color=RGBColor(0x88, 0x99, 0xBB),
             alignment=PP_ALIGN.CENTER)


# ============================================================
# 第2页：课题背景与意义
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '一、课题背景与意义', '为什么要做这个系统？')

# 左侧：行业现状
add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(2.8),
         '行业现状', [
             '▸ 健身俱乐部超4万家，健身人口突破7000万',
             '▸ 超过60%健身房无专业管理系统',
             '▸ 会员续费率不到20%',
             '▸ 本地部署成本几万起步',
         ], bg_color=ACCENT, font_size=15)

# 中间：现有问题
add_card(slide, Inches(4.7), Inches(1.5), Inches(3.8), Inches(2.8),
         '核心痛点', [
             '▸ AI建议千人一面，无个性化',
             '▸ 模型挂了→用户白屏报错',
             '▸ 多租户数据隔离缺失',
             '▸ 缺少订阅计费与激励机制',
         ], bg_color=RGBColor(0xDC, 0x35, 0x45), font_size=15)

# 右侧：解决方案
add_card(slide, Inches(8.9), Inches(1.5), Inches(3.8), Inches(2.8),
         '本系统目标', [
             '▸ SaaS多租户 + 行级隔离',
             '▸ AI上下文注入 → 因人而异',
             '▸ 多模型适配 + 降级容错',
             '▸ 三级订阅 + 成就激励闭环',
         ], bg_color=PRIMARY, font_size=15)

# 底部：政策背景和意义
add_shape(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.8), Inches(4.7), Inches(11.8), Inches(0.4),
             '研究意义', font_size=20, bold=True, color=PRIMARY_DARK)

meaning_items = [
    '▸ 政策驱动：《全民健身计划（2021-2025年）》提出经常参加体育锻炼人数比例达38.5%',
    '▸ SaaS模式降低门槛：中小健身房按月订阅，无需买服务器装软件',
    '▸ AI+数据闭环：训练记录 → 成就激励 → AI个性化建议，形成"记录-激励-优化"正循环',
    '▸ 填补空白：国内健身SaaS产品缺少真正的AI个性化 + 多模型降级容错能力',
]
add_bullet_list(slide, Inches(0.8), Inches(5.15), Inches(11.8), Inches(1.8),
                meaning_items, font_size=15, bold_prefix=False)

add_page_number(slide, 2, TOTAL_PAGES)


# ============================================================
# 第3页：课题内容 - 九大功能模块总览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '二、课题内容 — 九大功能模块', '系统覆盖九个核心功能域')

modules = [
    ('用户管理', ['注册/登录/JWT认证', '四级角色权限控制', '个人资料与头像'], PRIMARY),
    ('健身计划', ['初/中/高级模板', 'CRUD + 我的计划', '计划绑训练日志'], RGBColor(0x2E, 0x7D, 0x32)),
    ('训练日志', ['日志 + 组数明细', '事务原子写入', '统计聚合查询'], RGBColor(0x00, 0x89, 0x7B)),
    ('健康追踪', ['体重/体脂/心率/血压', '趋势折线图', '同日upsert去重'], RGBColor(0x55, 0x8B, 0x2F)),
    ('健身房管理', ['场馆信息维护', '会员入会/续费', 'gym_id行级隔离'], RGBColor(0xE6, 0x51, 0x00)),
    ('动作库', ['23个预置动作', '8大肌群分类', '动作与组数关联'], RGBColor(0xBF, 0x36, 0x0C)),
    ('成就勋章', ['17项勋章/4维度', '训练后自动检查', '进度条+解锁提示'], ACCENT),
    ('订阅计费', ['免费/月/年三级', '事务原子订阅变更', '过期自动标记'], RGBColor(0xAD, 0x14, 0x57)),
    ('AI健身助手', ['上下文注入个性化', '多模型适配切换', '降级容错机制'], RGBColor(0x8E, 0x24, 0xAA)),
]

cols = 3
rows = 3
card_w = Inches(3.9)
card_h = Inches(1.65)
gap_x = Inches(0.35)
gap_y = Inches(0.25)
start_x = Inches(0.5)
start_y = Inches(1.5)

for idx, (title, items, color) in enumerate(modules):
    r = idx // cols
    c = idx % cols
    x = start_x + c * (card_w + gap_x)
    y = start_y + r * (card_h + gap_y)

    # 标题栏
    add_shape(slide, x, y, card_w, Inches(0.45), color)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.05), card_w - Inches(0.2),
                 Inches(0.35), title, font_size=16, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    # 内容
    body = add_shape(slide, x, y + Inches(0.45), card_w, card_h - Inches(0.45), WHITE)
    body.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    for i, item in enumerate(items):
        add_text_box(slide, x + Inches(0.15), y + Inches(0.5) + i * Inches(0.35),
                     card_w - Inches(0.3), Inches(0.35),
                     item, font_size=13, color=BLACK)

add_page_number(slide, 3, TOTAL_PAGES)


# ============================================================
# 第4页：课题内容 - 核心创新点
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '二、课题内容 — 三大核心创新', '区别于现有健身SaaS的关键设计')

# 创新1：AI上下文注入
add_card(slide, Inches(0.5), Inches(1.5), Inches(3.9), Inches(5.3),
         '创新1：AI上下文注入个性化', [
             '问题：固定提示词 → 千人一面',
             '',
             '解决：实时构建用户上下文',
             '  • 个人信息（年龄/性别/目标）',
             '  • 近30天训练记录（10条）',
             '  • 最新5条健康指标',
             '  • 累计统计 + 连续打卡天数',
             '',
             '效果：不同用户获得差异化建议',
             '',
             '三类专项提示词：',
             '  • 训练建议（5维度引导）',
             '  • 营养指导（5维度展开）',
             '  • 计划推荐（5维度生成）',
         ], bg_color=RGBColor(0x8E, 0x24, 0xAA),
         content_bg=RGBColor(0xF3, 0xE5, 0xF5), font_size=13)

# 创新2：多模型适配+降级
add_card(slide, Inches(4.7), Inches(1.5), Inches(3.9), Inches(5.3),
         '创新2：多模型适配与降级容错', [
             '问题：单模型 → 一挂全瘫',
             '',
             '解决：统一接口+策略路由',
             '  • OpenAI / Anthropic / DeepSeek',
             '  • 本地 Ollama 兜底',
             '  • AI_PROVIDER环境变量切换',
             '',
             '降级机制：',
             '  • 密钥缺失 → 预设知识模板',
             '  • 服务不可达 → 领域FAQ回退',
             '  • is_fallback标识来源',
             '  • 预设回复50ms内返回',
             '',
             '速率限制：20次/用户/天',
         ], bg_color=RGBColor(0xE6, 0x51, 0x00),
         content_bg=RGBColor(0xFF, 0xF3, 0xE0), font_size=13)

# 创新3：SaaS多租户+订阅激励
add_card(slide, Inches(8.9), Inches(1.5), Inches(3.9), Inches(5.3),
         '创新3：SaaS多租户+订阅激励', [
             '问题：数据混在一起 / 缺商业化',
             '',
             '解决1：gym_id行级隔离',
             '  • 四级角色权限控制',
             '  • 前后端双重权限校验',
             '  • requireRole中间件',
             '',
             '解决2：三级订阅计费',
             '  • 免费/月29.9/年299',
             '  • 免费用户先用再转化',
             '  • 订阅变更事务原子保证',
             '',
             '解决3：成就激励闭环',
             '  • 17项勋章/4维度',
             '  • 训练完成→自动检查→即时解锁',
         ], bg_color=PRIMARY,
         content_bg=RGBColor(0xE8, 0xEE, 0xFA), font_size=13)

add_page_number(slide, 4, TOTAL_PAGES)


# ============================================================
# 第5页：毕设方案 - 系统架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '三、毕设方案 — 系统架构设计', '前后端分离三层架构')

# 表示层
add_shape(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(2), Inches(0.4),
             '表示层', font_size=18, bold=True, color=PRIMARY_DARK)
fe_items = [
    'React 19 + TypeScript + Ant Design 6 + Vite 8',
    'Zustand 5 状态管理（替代Redux，代码量减少60%）',
    'React Router v7 嵌套路由 + AuthGuard / PublicAuthGuard',
    'Axios apiClient统一请求：自动注入Bearer Token，401自动跳转登录',
]
add_bullet_list(slide, Inches(0.7), Inches(1.95), Inches(11.8), Inches(1.0),
                fe_items, font_size=14, bold_prefix=False)

# 业务逻辑层
add_shape(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(1.8), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(0.7), Inches(3.25), Inches(2.5), Inches(0.4),
             '业务逻辑层', font_size=18, bold=True, color=RGBColor(0x2E, 0x7D, 0x32))
be_items = [
    'Express 5：六层中间件栈（CORS白名单 → 请求体解析 → JWT认证 → 角色校验 → Zod校验 → 统一错误处理）',
    'Zod 4：一份Schema两端生效 — 运行时校验 + 自动推导TypeScript类型，12个校验Schema',
    'JWT + requireRole：认证层+授权层分离，支持多角色参数 requireRole("admin", "gym_admin")',
    'AI服务层：createChatCompletion统一接口，按AI_PROVIDER策略路由至不同模型端点',
]
add_bullet_list(slide, Inches(0.7), Inches(3.65), Inches(11.8), Inches(1.3),
                be_items, font_size=14, bold_prefix=False)

# 数据访问层
add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5), RGBColor(0xFF, 0xF3, 0xE0))
add_text_box(slide, Inches(0.7), Inches(5.25), Inches(2.5), Inches(0.4),
             '数据访问层', font_size=18, bold=True, color=RGBColor(0xE6, 0x51, 0x00))
da_items = [
    'PostgreSQL 16：ACID事务（训练日志+组数原子写入、订阅变更原子保证）',
    'pg.Pool连接池：最大20连接，空闲30s超时，2s连接超时',
    '原生SQL：四表LEFT JOIN + JSON聚合（比Prisma性能提升3倍），参数化查询防注入',
]
add_bullet_list(slide, Inches(0.7), Inches(5.65), Inches(11.8), Inches(1.0),
                da_items, font_size=14, bold_prefix=False)

# 右侧部署示意
add_shape(slide, Inches(9.0), Inches(6.85), Inches(3.8), Inches(0.5), LIGHT_GRAY)
add_text_box(slide, Inches(9.1), Inches(6.88), Inches(3.6), Inches(0.4),
             '部署：Docker Compose (postgres + backend + frontend)',
             font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 5, TOTAL_PAGES)


# ============================================================
# 第6页：毕设方案 - 数据模型设计
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '三、毕设方案 — 数据模型设计', '11张表，按业务域分组')

# 6个域名
domains = [
    ('用户域', 'users表\n• role: VARCHAR+CHECK\n  (user/coach/gym_admin/admin)\n• status: active/inactive/banned\n• email/username: UNIQUE', PRIMARY),
    ('健身房域', 'gyms + gym_members\n• 多对多关系\n• UNIQUE(gym_id,user_id)\n  防重复加入\n• membership_status区分状态', RGBColor(0xE6, 0x51, 0x00)),
    ('训练域', 'fitness_plans + workout_logs\n+ workout_sets + exercises\n• 组数拆为独立表（索引优化）\n• 复合索引(user_id,workout_date)\n  查询200ms→15ms', RGBColor(0x2E, 0x7D, 0x32)),
    ('健康数据域', 'health_data\n• UNIQUE(user_id,record_date)\n• INSERT ON CONFLICT UPDATE\n  (upsert防锯齿跳变)\n• 体脂/心率/血压/肌肉量', RGBColor(0x00, 0x89, 0x7B)),
    ('成就域', 'achievements + user_achievements\n• 17项勋章/4维度\n• code: 字符串可读ID\n• requirement_type:\n  workouts/days/duration/calories', ACCENT),
    ('订阅+AI域', 'subscriptions\n• 三级套餐/四种状态\n• 起止日期+金额\n• 事务原子订阅变更\n• 过期自动标记expired', RGBColor(0x8E, 0x24, 0xAA)),
]

card_w2 = Inches(3.9)
card_h2 = Inches(2.3)
gap2 = Inches(0.3)
for idx, (title, content, color) in enumerate(domains):
    r = idx // 3
    c = idx % 3
    x = Inches(0.5) + c * (card_w2 + gap2)
    y = Inches(1.5) + r * (card_h2 + Inches(0.2))

    add_shape(slide, x, y, card_w2, Inches(0.5), color)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.05), card_w2 - Inches(0.2),
                 Inches(0.4), title, font_size=16, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    body = add_shape(slide, x, y + Inches(0.5), card_w2, card_h2 - Inches(0.5), WHITE)
    body.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    # 手动拆行
    lines = content.split('\n')
    for i, line in enumerate(lines):
        add_text_box(slide, x + Inches(0.15), y + Inches(0.55) + i * Inches(0.25),
                     card_w2 - Inches(0.3), Inches(0.25),
                     line, font_size=12, color=BLACK)

# 底部索引设计要点
add_shape(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.8), LIGHT_GRAY)
add_text_box(slide, Inches(0.8), Inches(6.45), Inches(11.8), Inches(0.7),
             '索引设计原则：根据实际查询模式逐步添加 — (user_id, workout_date)复合索引使查询从200ms降至15ms；health_data同理',
             font_size=14, color=GRAY)

add_page_number(slide, 6, TOTAL_PAGES)


# ============================================================
# 第7页：毕设方案 - 关键实现细节
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '三、毕设方案 — 关键实现细节', '训练日志事务、AI上下文构建、降级容错')

# 左侧：训练日志事务
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.6), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(5.5), Inches(0.35),
             '训练日志：事务原子写入', font_size=18, bold=True, color=PRIMARY_DARK)
txn_items = [
    '▸ BEGIN → 插入workout_logs主表 → 逐条INSERT workout_sets从表 → COMMIT',
    '▸ 任一步骤失败 → ROLLBACK回滚，不留残数据',
    '▸ 事务完成后 → checkAndUnlockAchievements检查成就解锁',
    '▸ 更新策略：先删后插（DELETE + 重新INSERT完整sets数组）',
    '▸ 踩坑：初版组数用JSONB存 → 无法索引 → 拆独立表后性能提升10倍',
]
add_bullet_list(slide, Inches(0.7), Inches(1.9), Inches(5.5), Inches(2.0),
                txn_items, font_size=14, bold_prefix=False)

# 右侧：AI上下文构建
add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.6), RGBColor(0xF3, 0xE5, 0xF5))
add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5.5), Inches(0.35),
             'AI上下文构建流程', font_size=18, bold=True, color=RGBColor(0x8E, 0x24, 0xAA))
ctx_items = [
    '▸ 4个SQL查询并行：用户信息 + 30天训练 + 5条健康 + 累计统计',
    '▸ 连续打卡：SQL窗口函数高效计算',
    '▸ formatUserContextForPrompt → 结构化文本拼入系统提示词',
    '▸ 示例："累计训练:45次；连续打卡:7天；最近训练:上肢力量60min"',
    '▸ Token消耗增加约40%，但换来个性化区分度',
]
add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(2.0),
                ctx_items, font_size=14, bold_prefix=False)

# 下方：降级容错
add_shape(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8), RGBColor(0xFF, 0xF3, 0xE0))
add_text_box(slide, Inches(0.7), Inches(4.35), Inches(5), Inches(0.35),
             'AI降级容错机制', font_size=18, bold=True, color=RGBColor(0xE6, 0x51, 0x00))

# 流程图式描述
flow_items = [
    '请求到达 → isAIConfigured检查 → 密钥存在？',
    '    ✓ → createChatCompletion调用模型 → 成功？→ 返回AI回复（is_fallback: false）',
    '                                          → 失败？→ FALLBACK_RESPONSES预设回复（is_fallback: true）',
    '    ✗ → 直接返回FALLBACK_RESPONSES预设回复（is_fallback: true）',
    '',
    '前端标识：data.is_fallback === true → 橙色"预设建议"标签',
    '预设回复50ms内返回，保证核心功能不受第三方AI波动影响',
]
add_bullet_list(slide, Inches(0.7), Inches(4.75), Inches(11.8), Inches(2.2),
                flow_items, font_size=14, bold_prefix=False)

add_page_number(slide, 7, TOTAL_PAGES)


# ============================================================
# 第8页：毕设方案 - 安全与部署
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '三、毕设方案 — 安全设计与容器化部署', '前后端双重安全 + Docker一键部署')

# 安全设计
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.3), RGBColor(0xFC, 0xE4, 0xEC))
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(5.5), Inches(0.35),
             '安全设计', font_size=20, bold=True, color=RGBColor(0xAD, 0x14, 0x57))

security_items = [
    '认证安全：',
    '  • bcrypt哈希存储密码（saltRounds=10）',
    '  • JWT 7天有效期，Bearer方案传递',
    '  • JWT_SECRET启动时强制校验，缺失则process.exit(1)',
    '',
    '授权安全：',
    '  • 后端：authMiddleware(认证) + requireRole(授权)两层',
    '  • 前端：hasRole()动态隐藏菜单和按钮',
    '  • 即使绕过前端，后端403依然拦截',
    '',
    '输入安全：',
    '  • Zod Schema覆盖全部API输入（12个校验模式）',
    '  • 参数化查询杜绝SQL注入',
    '  • 文件上传限制5MB',
    '  • AI速率限制：20次/用户/天',
    '',
    '数据安全：',
    '  • CORS白名单限制请求来源',
    '  • gym_id行级隔离场馆数据（3张表）',
]
add_bullet_list(slide, Inches(0.7), Inches(1.95), Inches(5.5), Inches(4.6),
                security_items, font_size=13, bold_prefix=False)

# 部署
add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5.5), Inches(0.35),
             'Docker容器化部署', font_size=20, bold=True, color=RGBColor(0x2E, 0x7D, 0x32))

deploy_items = [
    'docker-compose.yml 三个服务：',
    '  • postgres: PostgreSQL 15 Alpine + 健康检查',
    '  • backend: Node.js多阶段构建 + HTTP探针',
    '  • frontend: Nginx托管静态资源 + API代理',
    '',
    '关键设计：',
    '  • depends_on + condition: service_healthy',
    '    → 确保DB就绪后再启动后端',
    '  • 后端Dockerfile多阶段构建',
    '    → 减小镜像体积',
    '  • 数据卷持久化DB和上传文件',
    '',
    '健康检查机制：',
    '  • PostgreSQL: pg_isready',
    '  • Backend: GET /health (12ms响应)',
    '',
    '扩展性：',
    '  • 单台2核4G云服务器承载数百并发',
    '  • 负载均衡器后加节点即可水平扩展',
]
add_bullet_list(slide, Inches(7.0), Inches(1.95), Inches(5.5), Inches(4.6),
                deploy_items, font_size=13, bold_prefix=False)

add_page_number(slide, 8, TOTAL_PAGES)


# ============================================================
# 第9页：毕设成果 - 测试结果
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, '四、毕设成果 — 测试验证', '功能测试 + 性能测试 + 降级验证')

# 性能数据卡片
add_shape(slide, Inches(0.5), Inches(1.5), Inches(8.0), Inches(2.5), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(7.5), Inches(0.35),
             '性能测试结果（autocannon, 100次请求/端点）', font_size=18, bold=True, color=PRIMARY_DARK)

perf_data = [
    '指标                              平均响应    P99',
    '─────────────────────────────────────────',
    'GET /health（健康检查）       12ms          -',
    'CRUD接口（最高频10个）        38-68ms      ≤180ms',
    'AI降级预设回复                   50ms          -',
]
for i, line in enumerate(perf_data):
    color = PRIMARY if i == 0 else (GRAY if i == 1 else BLACK)
    font_size = 14 if i > 1 else 15
    bold = i == 0
    add_text_box(slide, Inches(0.9), Inches(1.95) + i * Inches(0.3),
                 Inches(7.3), Inches(0.3), line, font_size=font_size,
                 color=color, bold=bold, font_name='Consolas')

# 右侧关键结论
add_shape(slide, Inches(8.8), Inches(1.5), Inches(4.0), Inches(2.5), LIGHT_GRAY)
add_text_box(slide, Inches(9.0), Inches(1.55), Inches(3.6), Inches(0.35),
             '关键结论', font_size=18, bold=True, color=PRIMARY_DARK)
conclusions = [
    '▸ CRUD接口全部满足200ms目标',
    '▸ AI降级100%触发，预设回复50ms返回',
    '▸ AI模块代码完整，降级模式稳定运行',
    '▸ 配置API Key后即可调用真实模型',
    '▸ 事务回滚测试：主表+从表同步回滚',
    '▸ Zod防御：注入类输入全部在中间件层拦截',
]
add_bullet_list(slide, Inches(9.0), Inches(1.95), Inches(3.6), Inches(1.8),
                conclusions, font_size=13, bold_prefix=False)

# 功能测试要点
add_shape(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.9), LIGHT_GRAY)
add_text_box(slide, Inches(0.7), Inches(4.25), Inches(11.8), Inches(0.35),
             '功能测试要点与踩坑经验', font_size=18, bold=True, color=PRIMARY_DARK)

test_items = [
    '▸ 用户名200+字符 → JWT膨胀到2KB → 间歇性401 → Zod限制3-20位',
    '▸ 主表插入成功但从表失败 → 残数据训练日志 → 加事务后ROLLBACK全部回滚',
    '▸ 健康数据同日两次录入 → 折线图锯齿跳变 → UNIQUE约束 + UPSERT语义',
    '▸ 免费用户访问AI页 → 后端返回403白屏 → 前端加套餐检查+升级提示',
    '▸ 月度升年度并发问题 → 取消+创建不在同一事务 → 包进BEGIN-COMMIT',
    '▸ 成就系统：训练达10次 → "小试身手"解锁，再次调用不重复解锁',
]
add_bullet_list(slide, Inches(0.7), Inches(4.65), Inches(11.8), Inches(2.3),
                test_items, font_size=14, bold_prefix=False)

add_page_number(slide, 9, TOTAL_PAGES)


# ============================================================
# 第10页：总结与致谢
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 大背景
add_shape(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, PRIMARY_DARK)
add_shape(slide, Inches(0), Inches(3.2), prs.slide_width, Inches(0.04), ACCENT)

add_text_box(slide, Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), Inches(0.6),
             '总结与致谢', font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# 工作总结
sum_items = [
    '▸ 覆盖9大功能模块、50+ API端点、11张数据表',
    '▸ 实现"记录-激励-优化"闭环：训练日志→成就激励→AI上下文',
    '▸ AI上下文注入让不同用户获得差异化建议，降级机制100%触发',
    '▸ CRUD接口平均38-68ms，P99≤180ms',
]
y = Inches(1.3)
for item in sum_items:
    add_text_box(slide, Inches(1.0), y, Inches(11.3), Inches(0.4),
                 item, font_size=18, color=WHITE)
    y += Inches(0.4)

# 未来工作
add_text_box(slide, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.4),
             '未来工作', font_size=22, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

future_items = [
    '▸ Schema级数据隔离（活跃场馆>100时升级）',
    '▸ AI响应流式化（SSE推送，首字延迟<1s）',
    '▸ Redis分布式速率计数 + WebSocket实时推送',
    '▸ 移动端深度适配 + 国际化i18n + 深色模式',
]
y = Inches(3.9)
for item in future_items:
    add_text_box(slide, Inches(1.0), y, Inches(11.3), Inches(0.4),
                 item, font_size=17, color=RGBColor(0xBB, 0xCC, 0xEE))
    y += Inches(0.38)

# 致谢
add_shape(slide, Inches(0), Inches(5.6), prs.slide_width, Inches(0.04), ACCENT)
add_text_box(slide, Inches(0.5), Inches(5.9), prs.slide_width - Inches(1), Inches(0.5),
             '感谢王京文老师的悉心指导', font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.5), Inches(6.3), prs.slide_width - Inches(1), Inches(0.5),
             '感谢各位评委老师的聆听与指导', font_size=20, color=RGBColor(0xBB, 0xCC, 0xEE),
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.5), Inches(6.8), prs.slide_width - Inches(1), Inches(0.5),
             '陈斌 / 202203010223 / 计算机2202', font_size=16, color=RGBColor(0x88, 0x99, 0xBB),
             alignment=PP_ALIGN.CENTER)

add_page_number(slide, 10, TOTAL_PAGES)


# ========== 保存 ==========
output_path = os.path.join(os.path.dirname(__file__), '毕业答辩PPT-基于SaaS模式的智慧健身管理系统.pptx')
prs.save(output_path)
print(f'PPT已生成：{output_path}')
print(f'共 {len(prs.slides)} 页')